from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml import parse_xml
from pptx.dml.line import LineFormat
from pptx import Presentation
from pptx.util import Pt
import math
import util.const as const
import sys
sys.path.append('../util')


class PPTXSaver():
    def __init__(self, filepath, painter):
        self.filepath = filepath
        self.fishbone = painter.fishbone

    def add_textbox(self, slide, x, y, w, h, halign, valign, font_size, text, color, is_bold=False):
        textbox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        text_frm = textbox.text_frame
        text_frm.paragraphs[0].alignment = halign
        text_frm.paragraphs[0].font.size = Pt(font_size)
        text_frm.paragraphs[0].font.color.rgb = color
        text_frm.paragraphs[0].font.bold = is_bold
        text_frm.paragraphs[0].font.name = 'Meiryo'
        text_frm.paragraphs[0].text = text
        text_frm.vertical_anchor = valign
        text_frm.word_wrap = True

        return textbox

    def add_slide_bone_text(self, bone, slide, side):
        font_size = 9
        font_color = RGBColor(0, 0, 0)

        if bone.direction == 'vertical':
            if side == 'up':
                self.add_textbox(slide, bone.bone[0].x - bone.get_width() / 2, bone.bone[0].y - const.VERTICAL_MARGIN + 5, bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.CENTER, MSO_ANCHOR.BOTTOM, font_size, bone.text, font_color)
            else:
                self.add_textbox(slide, bone.bone[0].x - bone.get_width() / 2, bone.bone[0].y + 10, bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.CENTER, MSO_ANCHOR.TOP, font_size, bone.text, font_color)
        else:
            if side == 'up':
                self.add_textbox(slide, bone.bone[0].x, bone.bone[0].y + 5, bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_size, bone.text, font_color)

            else:
                self.add_textbox(slide, bone.bone[0].x, bone.bone[0].y - const.VERTICAL_MARGIN, bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.LEFT, MSO_ANCHOR.BOTTOM, font_size, bone.text, font_color)

        for child_bone in bone.child_bones:
            self.add_slide_bone_text(child_bone, slide, side)

    def add_line(self, slide, x1, y1, x2, y2):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
        line_elem = connector.line._get_or_add_ln()
        line_elem.append(parse_xml("""
            <a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
        """))
        line = LineFormat(connector)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)

    def add_slide_bone_line(self, bone, slide):
        self.add_line(
            slide, bone.bone[0].x, bone.bone[0].y, bone.bone[1].x, bone.bone[1].y)

        for child_bone in bone.child_bones:
            self.add_slide_bone_line(child_bone, slide)

    def save(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        main_bone_arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Pt(self.fishbone.main_bone[0].x), Pt(
                self.fishbone.main_bone[0].y - 30),
            Pt(self.fishbone.main_bone[1].x - self.fishbone.main_bone[0].x),
            Pt(60))

        color = (int(const.MAIN_BONE_TEXT_COLOR[0] * 255), int(
            const.MAIN_BONE_TEXT_COLOR[1] * 255), int(const.MAIN_BONE_TEXT_COLOR[2] * 255))

        main_bone_text_color = RGBColor(
            color[0], color[1], color[2])

        self.add_textbox(
            slide,
            self.fishbone.main_bone[1].x, self.fishbone.main_bone[1].y - 100,
            200, 200, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 20, self.fishbone.text, main_bone_text_color, True)

        color = (int(const.MAIN_SUB_BONE_COLOR[0] * 255), int(
            const.MAIN_SUB_BONE_COLOR[1] * 255), int(const.MAIN_SUB_BONE_COLOR[2] * 255))

        main_sub_bone_color = RGBColor(
            color[0], color[1], color[2])

        main_bone_arrow.fill.solid()
        main_bone_arrow.fill.fore_color.rgb = main_sub_bone_color

        index = 0
        side = ''
        sub_bone_font_size = 20

        for sub_bone in self.fishbone.sub_bones:

            if index % 2 == 0:
                side = 'up'
                width = (sub_bone.bone[1].y - sub_bone.bone[0].y) / \
                    math.cos(math.radians(const.DEG)) - 15

                sub_bone_line = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Pt(sub_bone.bone[0].x), Pt(sub_bone.bone[0].y - 6),
                    Pt(width), Pt(12)
                )
                sub_bone_line.rotation = 90 - const.DEG

                dx = (width / 2) * \
                    math.sin(math.radians(const.DEG)) - (width / 2)
                dy = (width / 2) * math.cos(math.radians(const.DEG))

                sub_bone_line.left = sub_bone_line.left + Pt(dx)
                sub_bone_line.top = sub_bone_line.top + Pt(dy)

                self.add_textbox(slide, sub_bone.bone[0].x - sub_bone.get_width() / 2, sub_bone.bone[0].y - const.VERTICAL_MARGIN, sub_bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.CENTER, MSO_ANCHOR.BOTTOM, sub_bone_font_size, sub_bone.text, main_sub_bone_color)
            else:
                side = 'down'
                width = (sub_bone.bone[0].y - sub_bone.bone[1].y) / \
                    math.cos(math.radians(const.DEG)) - 15

                sub_bone_line = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Pt(sub_bone.bone[0].x), Pt(sub_bone.bone[0].y - 6),
                    Pt(width), Pt(12)
                )
                sub_bone_line.rotation = -(90 - const.DEG)

                dx = (width / 2) * \
                    math.sin(math.radians(const.DEG)) - (width / 2)
                dy = (width / 2) * math.cos(math.radians(const.DEG))

                sub_bone_line.left = sub_bone_line.left + Pt(dx)
                sub_bone_line.top = sub_bone_line.top - Pt(dy)

                self.add_textbox(slide, sub_bone.bone[0].x - sub_bone.get_width() / 2, sub_bone.bone[0].y + 20, sub_bone.get_width(
                ), const.VERTICAL_MARGIN, PP_ALIGN.CENTER, MSO_ANCHOR.TOP, sub_bone_font_size, sub_bone.text, main_sub_bone_color)

            index = index + 1

            sub_bone_line.fill.solid()
            sub_bone_line.fill.fore_color.rgb = RGBColor(
                color[0], color[1], color[2])

            for bone in sub_bone.child_bones:
                self.add_slide_bone_line(bone, slide)
                self.add_slide_bone_text(bone, slide, side)

        prs.save(self.filepath)
